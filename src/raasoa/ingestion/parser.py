"""Document parsing for multiple formats.

Extracts text, tables, and metadata from:
- PDF (pypdf with table detection)
- DOCX (python-docx with tables and styles)
- XLSX/CSV (rows → structured text with headers)
- PPTX (slides with notes)
- TXT/MD/HTML (plain text)

Tables are preserved as markdown-formatted text so the chunker
and claim extractor can understand tabular data.
"""

from __future__ import annotations

import csv
import io
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any


@dataclass
class ParsedSection:
    text: str
    title: str | None = None
    section_type: str = "text"  # 'text', 'table', 'code', 'heading'
    page_number: int | None = None  # PDF page, PPTX slide
    source_location: str | None = None  # e.g. "Page 5", "Slide 3", "Sheet: Revenue"


@dataclass
class ParsedDocument:
    title: str | None
    full_text: str
    sections: list[ParsedSection] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)


def _table_to_markdown(headers: list[str], rows: list[list[str]]) -> str:
    """Convert a table to markdown format for better chunking."""
    if not headers and not rows:
        return ""

    if not headers and rows:
        headers = [f"Col{i + 1}" for i in range(len(rows[0]))]

    lines = []
    lines.append("| " + " | ".join(headers) + " |")
    lines.append("| " + " | ".join("---" for _ in headers) + " |")
    for row in rows:
        # Pad row to match header count
        padded = row + [""] * (len(headers) - len(row))
        lines.append("| " + " | ".join(padded[:len(headers)]) + " |")

    return "\n".join(lines)


def parse_text(content: str, filename: str) -> ParsedDocument:
    """Parse plain text, markdown, or similar formats."""
    lines = content.strip().split("\n")
    title = lines[0].strip("# ").strip() if lines else filename
    return ParsedDocument(
        title=title,
        full_text=content,
        sections=[ParsedSection(text=content, title=title)],
        metadata={"format": "text", "filename": filename},
    )


def parse_html(content: str, filename: str) -> ParsedDocument:
    """Parse HTML — strip tags, keep text structure."""
    import re

    # Basic HTML tag removal preserving structure
    text = re.sub(r"<br\s*/?>", "\n", content, flags=re.IGNORECASE)
    text = re.sub(r"</(p|div|h[1-6]|li|tr)>", "\n", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", "", text)
    # Decode common entities
    for entity, char in [("&amp;", "&"), ("&lt;", "<"), ("&gt;", ">"),
                          ("&quot;", '"'), ("&nbsp;", " ")]:
        text = text.replace(entity, char)
    # Clean up whitespace
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return parse_text(text, filename)


def parse_pdf(data: bytes, filename: str) -> ParsedDocument:
    """Parse PDF with text extraction and basic table detection."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages: list[str] = []
        sections: list[ParsedSection] = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(text)
                sections.append(
                    ParsedSection(
                        text=text,
                        title=f"Page {i + 1}",
                        section_type="text",
                        page_number=i + 1,
                        source_location=f"Page {i + 1}",
                    )
                )

            # Try to extract tables (pypdf 4+ supports this)
            try:
                tables = page.extract_tables()  # type: ignore[attr-defined]
                if tables:
                    for t_idx, table in enumerate(tables):
                        if not table or not table[0]:
                            continue
                        headers = [str(c or "") for c in table[0]]
                        rows = [
                            [str(c or "") for c in row]
                            for row in table[1:]
                        ]
                        md_table = _table_to_markdown(headers, rows)
                        if md_table:
                            pages.append(md_table)
                            sections.append(
                                ParsedSection(
                                    text=md_table,
                                    title=f"Page {i + 1} Table {t_idx + 1}",
                                    section_type="table",
                                    page_number=i + 1,
                                    source_location=f"Page {i + 1}, Table {t_idx + 1}",
                                )
                            )
            except (AttributeError, Exception):
                pass  # Old pypdf without table support

        full_text = "\n\n".join(pages)
        title = filename
        if reader.metadata and reader.metadata.title:
            title = reader.metadata.title

        meta: dict[str, Any] = {
            "format": "pdf",
            "filename": filename,
            "pages": len(reader.pages),
        }
        if reader.metadata:
            if reader.metadata.author:
                meta["author"] = reader.metadata.author
            if reader.metadata.subject:
                meta["subject"] = reader.metadata.subject
            if reader.metadata.creation_date:
                meta["created"] = str(reader.metadata.creation_date)

        return ParsedDocument(
            title=title, full_text=full_text,
            sections=sections, metadata=meta,
        )
    except ImportError as err:
        msg = (
            "pypdf is required for PDF parsing. "
            "Install with: uv sync --extra parsing"
        )
        raise ImportError(msg) from err


def parse_docx(data: bytes, filename: str) -> ParsedDocument:
    """Parse DOCX with paragraphs, headings, and tables."""
    try:
        from docx import Document as DocxDocument

        doc = DocxDocument(io.BytesIO(data))
        parts: list[str] = []
        sections: list[ParsedSection] = []
        current_title: str | None = None

        # Paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else ""
            if style_name.startswith("Heading"):
                current_title = text
                sections.append(
                    ParsedSection(
                        text=text, title=text, section_type="heading",
                    )
                )
            else:
                sections.append(
                    ParsedSection(text=text, title=current_title)
                )
            parts.append(text)

        # Tables
        for t_idx, table in enumerate(doc.tables):
            headers: list[str] = []
            rows: list[list[str]] = []

            for r_idx, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                if r_idx == 0:
                    headers = cells
                else:
                    rows.append(cells)

            md_table = _table_to_markdown(headers, rows)
            if md_table:
                parts.append(md_table)
                sections.append(
                    ParsedSection(
                        text=md_table,
                        title=f"Table {t_idx + 1}",
                        section_type="table",
                    )
                )

        full_text = "\n\n".join(parts)
        title = parts[0] if parts else filename

        # Extract core properties
        meta: dict[str, Any] = {"format": "docx", "filename": filename}
        if doc.core_properties:
            cp = doc.core_properties
            if cp.author:
                meta["author"] = cp.author
            if cp.title:
                title = cp.title
            if cp.subject:
                meta["subject"] = cp.subject

        return ParsedDocument(
            title=title, full_text=full_text,
            sections=sections, metadata=meta,
        )
    except ImportError as err:
        msg = (
            "python-docx is required for DOCX parsing. "
            "Install with: uv sync --extra parsing"
        )
        raise ImportError(msg) from err


def parse_csv(data: bytes, filename: str) -> ParsedDocument:
    """Parse CSV/TSV as structured data.

    Each row becomes a readable text block with column headers,
    and the full table is also rendered as markdown.
    """
    text = data.decode("utf-8", errors="replace").strip()
    if not text:
        return ParsedDocument(
            title=filename, full_text="(empty CSV)",
            metadata={"format": "csv", "filename": filename},
        )

    import contextlib

    detected_dialect: type[csv.Dialect] | None = None
    with contextlib.suppress(csv.Error):
        detected_dialect = csv.Sniffer().sniff(text[:2048])

    reader = csv.reader(io.StringIO(text), detected_dialect or csv.excel)

    rows_raw = list(reader)
    if not rows_raw:
        return ParsedDocument(
            title=filename, full_text="(empty CSV)",
            metadata={"format": "csv", "filename": filename},
        )

    headers = rows_raw[0]
    data_rows = rows_raw[1:]

    # Markdown table
    md_table = _table_to_markdown(headers, data_rows[:200])

    # Also create row-by-row text for claim extraction
    row_texts: list[str] = []
    for row in data_rows[:500]:  # Cap at 500 rows
        pairs = [
            f"{h}: {v}" for h, v in zip(headers, row, strict=False) if v.strip()
        ]
        if pairs:
            row_texts.append("; ".join(pairs))

    full_text = f"# {filename}\n\n{md_table}"
    if row_texts:
        full_text += "\n\n## Row Details\n\n" + "\n".join(row_texts)

    sections = [
        ParsedSection(text=md_table, title="Data Table", section_type="table"),
    ]
    if row_texts:
        sections.append(
            ParsedSection(
                text="\n".join(row_texts),
                title="Row Details",
                section_type="text",
            )
        )

    return ParsedDocument(
        title=filename,
        full_text=full_text,
        sections=sections,
        metadata={
            "format": "csv",
            "filename": filename,
            "rows": len(data_rows),
            "columns": len(headers),
            "headers": headers,
        },
    )


def parse_xlsx(data: bytes, filename: str) -> ParsedDocument:
    """Parse Excel files (.xlsx) as structured data."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        all_parts: list[str] = []
        all_sections: list[ParsedSection] = []
        total_rows = 0

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append(cells)

            if not rows:
                continue

            headers = rows[0]
            data_rows = rows[1:]
            total_rows += len(data_rows)

            md_table = _table_to_markdown(headers, data_rows[:200])
            if md_table:
                all_parts.append(f"## {sheet_name}\n\n{md_table}")
                all_sections.append(
                    ParsedSection(
                        text=md_table,
                        title=sheet_name,
                        section_type="table",
                        source_location=f"Sheet: {sheet_name}",
                    )
                )

        wb.close()

        full_text = f"# {filename}\n\n" + "\n\n".join(all_parts)
        return ParsedDocument(
            title=filename,
            full_text=full_text,
            sections=all_sections,
            metadata={
                "format": "xlsx",
                "filename": filename,
                "sheets": wb.sheetnames,
                "total_rows": total_rows,
            },
        )
    except ImportError as err:
        msg = (
            "openpyxl is required for Excel parsing. "
            "Install with: uv sync --extra parsing"
        )
        raise ImportError(msg) from err


def parse_pptx(data: bytes, filename: str) -> ParsedDocument:
    """Parse PowerPoint files (.pptx) — slides + notes."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []
        sections: list[ParsedSection] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                if shape.has_table:
                    table = shape.table
                    headers = [
                        cell.text.strip() for cell in table.rows[0].cells
                    ]
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for row in table.rows[1:]
                    ]
                    md = _table_to_markdown(headers, rows)
                    if md:
                        slide_texts.append(md)

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Speaker Notes: {notes}]")

            if slide_texts:
                slide_text = "\n\n".join(slide_texts)
                parts.append(f"## Slide {i}\n\n{slide_text}")
                sections.append(
                    ParsedSection(
                        text=slide_text,
                        title=f"Slide {i}",
                        section_type="text",
                        page_number=i,
                        source_location=f"Slide {i}",
                    )
                )

        full_text = f"# {filename}\n\n" + "\n\n".join(parts)
        return ParsedDocument(
            title=filename,
            full_text=full_text,
            sections=sections,
            metadata={
                "format": "pptx",
                "filename": filename,
                "slides": len(prs.slides),
            },
        )
    except ImportError as err:
        msg = (
            "python-pptx is required for PowerPoint parsing. "
            "Install with: uv sync --extra parsing"
        )
        raise ImportError(msg) from err


def parse_file(data: bytes, filename: str) -> ParsedDocument:
    """Parse a file based on its extension."""
    suffix = Path(filename).suffix.lower()

    if suffix == ".pdf":
        return parse_pdf(data, filename)
    elif suffix in (".docx", ".doc"):
        return parse_docx(data, filename)
    elif suffix in (".xlsx", ".xls"):
        return parse_xlsx(data, filename)
    elif suffix in (".pptx", ".ppt"):
        return parse_pptx(data, filename)
    elif suffix == ".csv" or suffix == ".tsv":
        return parse_csv(data, filename)
    elif suffix == ".html":
        return parse_html(
            data.decode("utf-8", errors="replace"), filename,
        )
    else:
        return parse_text(
            data.decode("utf-8", errors="replace"), filename,
        )
